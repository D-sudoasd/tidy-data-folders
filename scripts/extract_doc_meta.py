#!/usr/bin/env python3
"""Extract lightweight identity metadata from document-class files.

v0.2: PDF + DOCX (+ optional PPTX props). Full SHA-256 for identity;
quick_fingerprint only for candidate screening. text_snip off by default.
"""
from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
import tempfile
from datetime import datetime, timezone
from pathlib import Path

# Formats with extractors
EXTRACT_EXTS = {".pdf", ".docx", ".pptx"}
# Listed but not auto-renamed without extract
LEGACY_EXTS = {".doc", ".ppt"}
OPTIONAL_EXTS = {".xlsx", ".xlsm"}
DOC_EXTS = EXTRACT_EXTS | LEGACY_EXTS | OPTIONAL_EXTS

DOI_RE = re.compile(r"\b10\.\d{4,9}/[-._;()/:A-Z0-9]+", re.IGNORECASE)
PUB_ID_RE = re.compile(
    r"(1-s2\.0-|s\d{4,}-|metals-\d{2}-|main(\s*\(\d+\))?\.pdf$)",
    re.IGNORECASE,
)


def clean_ws(s: str) -> str:
    return re.sub(r"\s+", " ", (s or "").replace("\x00", " ")).strip()


def quick_fingerprint(path: Path, max_bytes: int = 2 * 1024 * 1024) -> str:
    h = hashlib.sha256()
    size = path.stat().st_size
    with path.open("rb") as f:
        h.update(f.read(max_bytes))
        h.update(b"|")
        h.update(str(size).encode("ascii"))
    return h.hexdigest()


def sha256_full(path: Path, chunk: int = 1024 * 1024) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        while True:
            b = f.read(chunk)
            if not b:
                break
            h.update(b)
    return h.hexdigest()


def extract_pdf(path: Path, include_snip: bool) -> dict:
    out = {
        "title": "",
        "authors": [],
        "year": None,
        "venue": "",
        "doi": "",
        "page_count": None,
        "text_snip": "",
        "is_likely_scanned": False,
        "extract_engine": "none",
        "error": "",
    }
    try:
        from pypdf import PdfReader
    except ImportError:
        out["error"] = "pypdf not installed"
        return out

    try:
        reader = PdfReader(str(path), strict=False)
        out["page_count"] = len(reader.pages)
        meta = reader.metadata
        if meta:
            if meta.title:
                out["title"] = clean_ws(str(meta.title))
            if meta.author:
                authors = re.split(r"[,;]| and ", str(meta.author))
                out["authors"] = [clean_ws(a) for a in authors if clean_ws(a)]

        texts = []
        for page in reader.pages[:2]:
            try:
                t = page.extract_text() or ""
            except Exception:
                t = ""
            texts.append(t)
        full = clean_ws("\n".join(texts))
        if include_snip:
            out["text_snip"] = full[:2500]
        out["extract_engine"] = "pypdf"
        if len(full) < 40 and (out["page_count"] or 0) >= 1:
            out["is_likely_scanned"] = True

        m = DOI_RE.search(full)
        if not m and meta:
            blob = " ".join(
                str(x) for x in (meta.title, meta.subject, meta.creator) if x
            )
            m = DOI_RE.search(blob)
        if m:
            out["doi"] = m.group(0).rstrip(".)],;")

        years = re.findall(r"\b((?:19|20)\d{2})\b", full[:1200])
        for y in years:
            yi = int(y)
            if 1980 <= yi <= 2035:
                out["year"] = yi
                break

        bad_title = re.compile(
            r"^(abstract|keywords|introduction|doi:|www\.|journal pre-proof|article in press|"
            r"available online|contents lists|elsevier|springer|open access|research article|"
            r"regular article|accepted manuscript|proofs?|pii:|microsoft (word|powerpoint)|"
            r"page \d+|vol\.|pp\.|copyright|all rights reserved|providing this early version|"
            r"article are included|creative commons|graduate institute|department of|"
            r"university of|national laboratory|school of|college of)",
            re.I,
        )
        bad_title_contains = re.compile(
            r"creative commons|all rights reserved|corresponding author|e-?mail:|"
            r"https?://|journal pre-proof|article in press",
            re.I,
        )
        if out["title"] and bad_title.match(out["title"].strip()):
            out["title"] = ""
        if not out["title"] and full:
            lines = [
                clean_ws(L)
                for L in re.split(r"[\n\r]+", "\n".join(texts))
                if clean_ws(L)
            ]
            ranked = []
            for i, L in enumerate(lines[:25]):
                if len(L) < 15 or len(L) > 220:
                    continue
                if bad_title.match(L) or bad_title_contains.search(L):
                    continue
                if re.search(
                    r"\b(Institute|University|Laboratory|Department|School of|College of)\b",
                    L,
                ) and len(L) < 120:
                    continue
                if re.match(r"^\d+$", L):
                    continue
                if DOI_RE.search(L) and len(L) < 40:
                    continue
                letters = len(re.findall(r"[A-Za-z\u4e00-\u9fff]", L))
                if letters < 12:
                    continue
                ranked.append((1 if letters >= 30 else 0, letters, -i, L))
            if ranked:
                ranked.sort(reverse=True)
                out["title"] = ranked[0][3]

        # keep a short signal for classification even without snip persistence
        out["_class_snip"] = full[:800]
    except Exception as e:
        out["error"] = f"{type(e).__name__}: {e}"
    return out


def extract_docx(path: Path, include_snip: bool) -> dict:
    out = {
        "title": "",
        "authors": [],
        "year": None,
        "venue": "",
        "doi": "",
        "page_count": None,
        "text_snip": "",
        "is_likely_scanned": False,
        "extract_engine": "none",
        "error": "",
        "_class_snip": "",
    }
    try:
        from docx import Document  # type: ignore
    except ImportError:
        out["error"] = "python-docx not installed"
        return out
    try:
        doc = Document(str(path))
        props = doc.core_properties
        if props.title:
            out["title"] = clean_ws(props.title)
        if props.author:
            out["authors"] = [clean_ws(props.author)]
        paras = [clean_ws(p.text) for p in doc.paragraphs[:40] if clean_ws(p.text)]
        blob = clean_ws(" ".join(paras))
        out["_class_snip"] = blob[:800]
        if include_snip:
            out["text_snip"] = blob[:2500]
        out["extract_engine"] = "python-docx"
        if not out["title"] and paras:
            out["title"] = paras[0][:200]
        m = DOI_RE.search(blob)
        if m:
            out["doi"] = m.group(0).rstrip(".)],;")
        years = re.findall(r"\b((?:19|20)\d{2})\b", blob[:800])
        for y in years:
            yi = int(y)
            if 1980 <= yi <= 2035:
                out["year"] = yi
                break
    except Exception as e:
        out["error"] = f"{type(e).__name__}: {e}"
    return out


def extract_pptx(path: Path, include_snip: bool) -> dict:
    out = {
        "title": "",
        "authors": [],
        "year": None,
        "venue": "",
        "doi": "",
        "page_count": None,
        "text_snip": "",
        "is_likely_scanned": False,
        "extract_engine": "none",
        "error": "",
        "_class_snip": "",
    }
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        out["error"] = "python-pptx not installed"
        out["extract_engine"] = "none"
        return out
    try:
        prs = Presentation(str(path))
        out["page_count"] = len(prs.slides)
        core = prs.core_properties
        if core.title:
            out["title"] = clean_ws(core.title)
        if core.author:
            out["authors"] = [clean_ws(core.author)]
        texts = []
        if prs.slides:
            for shape in prs.slides[0].shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(clean_ws(shape.text))
        blob = clean_ws(" ".join(texts))
        out["_class_snip"] = blob[:800]
        if include_snip:
            out["text_snip"] = blob[:2500]
        if not out["title"] and texts:
            out["title"] = texts[0][:200]
        out["extract_engine"] = "python-pptx"
    except Exception as e:
        out["error"] = f"{type(e).__name__}: {e}"
    return out


def filename_signals(name: str) -> dict:
    stem = Path(name).stem
    return {
        "filename_has_publisher_id": bool(PUB_ID_RE.search(name)),
        "filename_looks_academic": bool(
            re.match(r"^(19|20)\d{2}_[A-Za-z]", stem)
            or (" - " in stem and re.search(r"(19|20)\d{2}", stem))
        ),
        "filename_year": _year_from_name(stem),
    }


def _year_from_name(stem: str):
    m = re.search(r"\b((?:19|20)\d{2})\b", stem)
    if m:
        y = int(m.group(1))
        if 1980 <= y <= 2035:
            return y
    return None


def process_file(path: Path, root: Path, include_snip: bool) -> dict:
    ext = path.suffix.lower()
    try:
        rel = str(path.resolve().relative_to(root.resolve())).replace("\\", "/")
    except ValueError:
        rel = path.name

    st = path.stat()
    mtime_utc = datetime.fromtimestamp(st.st_mtime, tz=timezone.utc).strftime(
        "%Y-%m-%dT%H:%M:%SZ"
    )

    rec = {
        "path": rel,  # relative only (F-18)
        "rel": rel,
        "name": path.name,
        "ext": ext,
        "bytes": st.st_size,
        "mtime_utc": mtime_utc,
        "quick_fingerprint": quick_fingerprint(path),
        "sha256_full": sha256_full(path),
        "content_hash": "",  # alias for full — used by propose
        "content_hash_short": "",
        "title": "",
        "authors": [],
        "year": None,
        "year_source": "",
        "venue": "",
        "doi": "",
        "page_count": None,
        "text_snip": "",
        "class_snip": "",
        "is_likely_scanned": False,
        "extract_engine": "none",
        "extract_supported": ext in EXTRACT_EXTS,
        "error": "",
        "filename_signals": filename_signals(path.name),
    }
    rec["content_hash"] = rec["sha256_full"]
    rec["content_hash_short"] = rec["sha256_full"][:8]

    if ext == ".pdf":
        got = extract_pdf(path, include_snip)
    elif ext == ".docx":
        got = extract_docx(path, include_snip)
    elif ext == ".pptx":
        got = extract_pptx(path, include_snip)
    elif ext in LEGACY_EXTS:
        got = {
            "error": f"legacy format {ext}: no auto-rename (manual review)",
            "extract_engine": "none",
        }
    elif ext in OPTIONAL_EXTS:
        got = {
            "error": f"{ext}: extract not auto-rename by default (review only)",
            "extract_engine": "none",
        }
    else:
        got = {"error": f"unsupported {ext}", "extract_engine": "none"}

    for k in (
        "title",
        "authors",
        "year",
        "venue",
        "doi",
        "page_count",
        "text_snip",
        "is_likely_scanned",
        "extract_engine",
        "error",
    ):
        if k in got and got[k] not in (None, "", []):
            rec[k] = got[k]
        elif k in got and rec.get(k) in (None, "", []):
            rec[k] = got[k]
    # Classification may use _class_snip in-memory only; never persist raw text by default.
    # Boolean / code signals only unless --include-class-snip.
    _cs = got.get("_class_snip") or ""
    rec["class_snip"] = ""  # filled only if caller requests include-class-snip later
    rec["has_abstract"] = bool(re.search(r"\babstract\b|摘要", _cs, re.I))
    rec["has_references"] = bool(re.search(r"\breferences\b|参考文献", _cs, re.I))
    rec["has_invoice_label"] = bool(re.search(r"发票号码|invoice\s*(no|number)", _cs, re.I))
    rec["_private_class_snip"] = _cs  # stripped before JSON write unless include-class-snip

    if rec["year"] is not None:
        rec["year_source"] = "body_or_meta"
    fy = rec["filename_signals"].get("filename_year")
    if fy and (rec["year"] is None or rec["year_source"] == "body_or_meta"):
        # filename preferred when present (F-12) — mark both
        if rec["year"] is None:
            rec["year"] = fy
            rec["year_source"] = "filename"
        else:
            rec["year_filename"] = fy

    if " - " in path.stem and not rec["venue"]:
        parts = [p.strip() for p in path.stem.split(" - ")]
        if len(parts) >= 3 and re.fullmatch(r"(19|20)\d{2}", parts[1] or ""):
            rec["venue"] = parts[0]
            if not rec["year"]:
                rec["year"] = int(parts[1])
                rec["year_source"] = "filename"
            if not rec["authors"] and len(parts) >= 3:
                rec["authors"] = [parts[2].split()[0]]
            if not rec["title"] and len(parts) >= 4:
                rec["title"] = " - ".join(parts[3:])

    # never persist absolute path
    return rec


def _is_symlink_or_reparse(path: Path) -> bool:
    try:
        if path.is_symlink():
            return True
        import os as _os

        if _os.name == "nt":
            import ctypes

            GetFileAttributesW = ctypes.windll.kernel32.GetFileAttributesW  # type: ignore
            attrs = GetFileAttributesW(str(path))
            if attrs != 0xFFFFFFFF and (attrs & 0x400):
                return True
    except OSError:
        return True
    return False


def iter_docs(root: Path, recursive: bool):
    if recursive:
        paths = sorted(
            p for p in root.rglob("*") if p.is_file() and p.suffix.lower() in DOC_EXTS
        )
    else:
        paths = sorted(
            p for p in root.iterdir() if p.is_file() and p.suffix.lower() in DOC_EXTS
        )
    for p in paths:
        if p.name.startswith("DOC_RENAME") or p.name.startswith("_doc_meta"):
            continue
        parts = {x.lower() for x in p.parts}
        if ".git" in parts or "__pycache__" in parts:
            continue
        # refuse symlinks / reparse points anywhere in the path under root
        try:
            cur = p
            while True:
                if _is_symlink_or_reparse(cur):
                    break  # skip this file
                if cur == root or cur.parent == cur:
                    yield p
                    break
                cur = cur.parent
                try:
                    cur.relative_to(root)
                except ValueError:
                    break
        except OSError:
            continue


def main(argv=None) -> int:
    ap = argparse.ArgumentParser(description="Extract document identity metadata")
    ap.add_argument("--root", required=True, help="Folder root")
    ap.add_argument(
        "--out",
        default="",
        help="Output JSONL (default: temp file; print path)",
    )
    ap.add_argument("--recursive", action="store_true", default=True)
    ap.add_argument("--no-recursive", action="store_true")
    ap.add_argument("--file", default="", help="Single file path under root")
    ap.add_argument(
        "--include-text-snip",
        action="store_true",
        help="Persist text_snip in JSONL (privacy-sensitive; off by default)",
    )
    ap.add_argument(
        "--include-class-snip",
        action="store_true",
        help="Persist class_snip free text (privacy-sensitive; off by default)",
    )
    ap.add_argument(
        "--overwrite-generated-output",
        action="store_true",
        help="Allow overwriting existing JSONL output (never source docs)",
    )
    ap.add_argument(
        "--max-file-mb",
        type=float,
        default=200.0,
        help="Skip files larger than this many MiB (default 200)",
    )
    args = ap.parse_args(argv)

    root = Path(args.root).resolve()
    if not root.is_dir():
        print(f"Root not found: {root}", file=sys.stderr)
        return 2

    import os

    if args.out:
        out_path = Path(args.out)
        if out_path.suffix.lower() != ".jsonl":
            print("ERROR: --out must end with .jsonl", file=sys.stderr)
            return 2
        if out_path.exists() and not args.overwrite_generated_output:
            print(f"ERROR: output exists (refuse overwrite): {out_path}", file=sys.stderr)
            return 2
        # refuse overwriting a source document under root
        try:
            res = out_path.resolve()
            res.relative_to(root)
            if res.suffix.lower() in DOC_EXTS:
                print(f"ERROR: --out collides with document under root: {out_path}", file=sys.stderr)
                return 2
            if res.is_file() and res.suffix.lower() in DOC_EXTS:
                print(f"ERROR: --out would overwrite a document: {out_path}", file=sys.stderr)
                return 2
        except ValueError:
            pass
        # if path exists as a non-jsonl document by content check of extension already done
        if out_path.exists() and out_path.suffix.lower() in DOC_EXTS:
            print(f"ERROR: --out collides with document: {out_path}", file=sys.stderr)
            return 2
    else:
        fd, tmp = tempfile.mkstemp(prefix="doc_meta_", suffix=".jsonl")
        os.close(fd)
        out_path = Path(tmp)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    recursive = not args.no_recursive
    include_snip = args.include_text_snip
    max_bytes = int(args.max_file_mb * 1024 * 1024)

    if args.file:
        fp = Path(args.file)
        if not fp.is_absolute():
            fp = root / fp
        fp = fp.resolve()
        try:
            fp.relative_to(root)
        except ValueError:
            print(f"file not under root: {fp}", file=sys.stderr)
            return 2
        if _is_symlink_or_reparse(fp):
            print(f"file is symlink/reparse (refused): {fp}", file=sys.stderr)
            return 2
        files = [fp]
    else:
        files = list(iter_docs(root, recursive=recursive))

    n = 0
    errors = 0
    lines: list[str] = []
    for i, p in enumerate(files, 1):
        try:
            # require relative path under root (no basename fallback)
            rel = str(p.resolve().relative_to(root)).replace("\\", "/")
            if p.stat().st_size > max_bytes:
                rec = {
                    "rel": rel,
                    "name": p.name,
                    "error": f"skipped: larger than {args.max_file_mb} MiB",
                    "extract_supported": False,
                }
                errors += 1
            else:
                rec = process_file(p, root, include_snip)
                rec["rel"] = rel
                if not include_snip:
                    rec["text_snip"] = ""
                if args.include_class_snip:
                    rec["class_snip"] = rec.pop("_private_class_snip", "") or ""
                else:
                    rec.pop("_private_class_snip", None)
                    rec["class_snip"] = ""
                if rec.get("error"):
                    errors += 1
        except ValueError as e:
            rec = {
                "rel": p.name,
                "name": p.name,
                "error": f"path_error: {e}",
            }
            errors += 1
        except Exception as e:
            rec = {
                "rel": p.name,
                "name": p.name,
                "error": f"{type(e).__name__}: {e}",
            }
            errors += 1
        rec.pop("abs_path", None)
        rec.pop("_private_class_snip", None)
        lines.append(json.dumps(rec, ensure_ascii=False))
        n += 1
        if i % 25 == 0:
            print(f"progress: {i}/{len(files)} errors={errors}", file=sys.stderr)

    text = "\n".join(lines) + ("\n" if lines else "")
    tmp = out_path.with_name(out_path.name + f".tmp.{os.getpid()}")
    try:
        with tmp.open("w", encoding="utf-8", newline="\n") as f:
            f.write(text)
            f.flush()
            os.fsync(f.fileno())
        os.replace(tmp, out_path)
    finally:
        if tmp.exists():
            try:
                tmp.unlink()
            except OSError:
                pass

    print(f"wrote {n} records -> {out_path}")
    print(f"errors: {errors}")
    print(
        "capabilities: pdf=pypdf docx=python-docx pptx=python-pptx(optional) "
        "doc/ppt=legacy-no-auto xlsx=review-only"
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
